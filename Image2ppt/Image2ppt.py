import os, os.path
from pptx import Presentation
from pptx.util import Inches
import tkinter
from tkinter import filedialog
from tkinter import *
from tkinter import ttk
# Utilizes python-pptx: https://python-pptx.readthedocs.io/

def main():

    create_gui = CreateGUI()
    append_ppt = AppendPPT()
    create_gui.construct_form()

    append_ppt.whole_process()


class CreateGUI:
    def __init__(self):
        self.root = Tk()
        self.root.minsize(width=500, height=300)
        self.root.title("GUI")
        self.frame = None
        self.label1 = None

    def construct_form(self):
        self.create_widget()
        label1 = self.create_label()
        button1 = self.create_button("Input path",self.get_path,0,1)

        button2 = self.create_button("Start",self.whole_process,1,1)

        self.start_widget()

    def create_widget(self):
        self.frame = ttk.Frame(self.root, padding=40)
        self.frame.grid()

    def create_label(self):
        label1 = ttk.Label(
            self.frame,
            text='Hello',
            background='#0000aa',
            foreground='#ffffff',
            padding=(5, 10))
        label1.grid(row=0, column=0)
        return label1

    def create_button(self,text,command,row,column):
        button1 = ttk.Button(
            self.frame,
            text=text,
            )
        button1.bind("<ButtonPress>",command)
        button1.grid(row=row, column=column)
        return button1

    def get_path(self, event):
        selected_path = filedialog.askdirectory()
        #self.event.widget.configure(text=selected_path)

    def whole_process(self,event):
        tkinter.Tk().withdraw()
        append_ppt = AppendPPT()
        input_path = append_ppt.get_path()
        output_path = append_ppt.get_path()

        append_slide = AppendSlide(input_path)

        prs = append_slide.append_images_in_ppt()
        prs.save(output_path + '/test.pptx')

    def start_widget(self):
        self.root.mainloop()


class AppendPPT:

    def get_path(self):
        selected_path = filedialog.askdirectory()
        return selected_path

    def whole_process(self):
        tkinter.Tk().withdraw()

        input_path = self.get_path()
        output_path = self.get_path()

        append_slide = AppendSlide(input_path)

        prs = append_slide.append_images_in_ppt()
        prs.save(output_path + '/test.pptx')


class AppendSlide:
    # initial value loading
    def __init__(self, input_path):
        self.column = 4
        self.row = 2

        self.ppt_width = 13.333
        self.ppt_height = 7.5

        self.img_iter = self.column * self.row

        self.img = self.get_images(input_path)
        self.img_count = len(self.img)

        self.input_path = input_path

    def get_images(self, input_path):
        img = os.listdir(input_path)
        return img

    # generate ppt and add images to the ppt
    def append_images_in_ppt(self):
        prs = Presentation()
        prs.slide_width = Inches(self.ppt_width)
        prs.slide_height = Inches(self.ppt_height)
        prs = self.append_images(prs)
        return prs

    def append_images(self, prs):
        blank_slide = prs.slide_layouts[6]

        for i in range(self.img_count):
            if i % self.img_iter == 0:
                image_slide = prs.slides.add_slide(blank_slide)

            arg1 = Inches((i % self.column) * (self.ppt_width / self.column))
            arg2 = Inches((i % self.img_iter // self.column) * self.ppt_height / self.row)

            image_slide.shapes.add_picture(self.input_path + '/' + self.img[i], arg1, arg2,
                                           height=Inches(self.ppt_width / self.column))

        return prs


if __name__ == "__main__":
    main()
