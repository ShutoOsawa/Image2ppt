import os, os.path
from pptx import Presentation
from pptx.util import Inches, Pt
import tkinter
from tkinter import filedialog
from tkinter import *
from tkinter import ttk
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
from math import ceil
import sys
import json
import math
import time
import Model
import View
import io

class Controller():
    def __init__(self):
        self.root = Tk()
        self.model = Model.Model()
        self.view = View.View(self.root)
        self.config = ConfigObject()
        if os.path.isfile("config.json"):
            with open('config.json', 'r') as f:
                data = json.load(f)
            print(data)
            print(type(data))
            self.config = LoadingConfig(data)

        self.view.input_path_label.configure(text=self.config.input_path)
        self.view.output_path_label.configure(text=self.config.output_path)

        self.bindings()
        self.ppt_component = SlideComponents()

        #self.config.input_path = self.view.input_path_label.cget("text")
        #self.config.output_path = self.view.output_path_label.cget("text")




    def bindings(self):
        self.view.input_path_button.bind("<ButtonPress>",
                                         lambda event: self.get_path(event, self.view.input_path_label))
        self.view.output_path_button.bind("<ButtonPress>",
                                          lambda event: self.get_path(event, self.view.output_path_label))
        self.view.start_process_button.bind("<ButtonPress>", lambda event: self.ppt_generation_process(event))
        self.view.save_config_button.bind("<ButtonPress>",
                                          lambda event: self.save_config_into_file(event, self.config))

    def save_config_into_file(self,event,arg):
        #s = json.dumps(arg, default=lambda x: x.__dict__)
        arg.input_path = self.view.input_path_label.cget("text")
        arg.output_path = self.view.output_path_label.cget("text")
        with open('config.json', 'w', encoding='utf-8') as f:
            json.dump(arg, f,default=lambda x: x.__dict__, ensure_ascii=False, indent=4)
        #print(s)

    def get_path(self, event, arg):
        selected_path = filedialog.askdirectory()
        if not selected_path=="":
            arg.configure(text=selected_path)

    def ppt_generation_process(self, event):
        #tkinter.Tk().withdraw()
        self.input_path = self.view.input_path_label.cget("text")
        self.output_path = self.view.input_path_label.cget("text")
        parameters = [self.view.gui_column.get(), self.view.gui_row.get(), self.view.gui_slide_counter.get(),
                      self.view.gui_ppt_width.get(), self.view.gui_ppt_height.get(),self.view.combobox.get()]

        self.get_parameters(parameters)
        prs = self.append_images_in_ppt()
        prs.save(self.output_path + '/' + self.view.gui_ppt_name_textbox.get() + '.pptx')
        os.startfile(self.output_path + '/' + self.view.gui_ppt_name_textbox.get() + '.pptx')


    def run(self):
        self.root.minsize(width=500, height=300)
        self.root.title("Image2ppt")
        self.root.mainloop()



    def get_parameters(self, parameters):
        self.combobox_value = parameters[5]
        self.img_list = self.get_images(self.input_path)
        self.img_count = len(self.img_list)
        #self.slide_number = 1

        self.column = int(parameters[0])
        self.row = int(parameters[1])
        self.ppt_width = float(parameters[3])
        self.ppt_height = float(parameters[4])
        self.img_width = self.ppt_width / self.column
        self.img_height = self.ppt_height / self.row
        self.img_iter = self.column * self.row
        self.slide_counter = int(parameters[2]) / self.img_iter

        print(parameters)
        print(self.combobox_value)


    def sort_images(self,list_of_files,dir_name):
        list_of_files = sorted(list_of_files,
                               key=lambda x: os.path.getmtime(os.path.join(dir_name, x))
                               , reverse=True
                               )
        return list_of_files

    def get_list_of_files(self,dir_name):
        list_of_files = filter(lambda x: os.path.isfile(os.path.join(dir_name, x)),
               os.listdir(dir_name))
        return list_of_files

    def get_images(self, input_path):
        # folder_files = os.listdir(input_path)
        img_list = []
        list_of_files = self.get_list_of_files(input_path)
        if self.combobox_value == "Ascending":
            list_of_files= self.sort_images(list_of_files,input_path)

        for file_name in list_of_files:
            self.model.check_image_extension(img_list,file_name)

        return img_list



    # generate ppt and add images to the ppt
    def append_images_in_ppt(self):
        prs = Presentation()
        prs.slide_width = Inches(self.ppt_width)
        prs.slide_height = Inches(self.ppt_height)
        prs = self.append_images(prs)
        return prs

    def append_images(self, prs):

        self.img_list = self.get_images(self.input_path)
        self.img_count = len(self.img_list)


        blank_slide = prs.slide_layouts[6]

        #get ppt width in pixels
        dpi = 72
        ppt_width_in_pixel = self.ppt_width*dpi
        ppt_height_in_pixel = self.ppt_height*dpi

        emus_per_px = int(914400 / dpi)
        # panel size in terms of pixel
        pixel_width = int(ppt_width_in_pixel / self.column)
        pixel_height = int(ppt_height_in_pixel / self.row)

        # panel width and height in inches
        #width = self.ppt_width / self.column
        #height = self.ppt_height / self.row

        #start adding images on the slide
        for i in range(self.img_count):
            #prepare blank slide if the image reaches threshold
            if i % self.img_iter == 0:
                image_slide = prs.slides.add_slide(blank_slide)
                #cellnumber = str(ceil(self.slide_number / self.slide_counter))
                cellnumber = str(ceil(len(prs.slides)/self.slide_counter))
                self.ppt_component.textbox(image_slide,cellnumber,self.ppt_width,self.ppt_height)
                #self.slide_number += 1
                #print(self.slide_number)
                print(len(prs.slides))

            #prepare image
            current_img = Image.open(self.input_path + '/' + self.img_list[i])

            #resize image
            ratio = self.model.get_resize_ratio(current_img.width, current_img.height, pixel_width, pixel_height)
            resized_img = current_img.resize((int(current_img.width * ratio), int(current_img.height * ratio)))

            #prepare margin for resized image
            margin_width = self.model.get_margin_in_pixel(pixel_width,resized_img.width)
            margin_height = self.model.get_margin_in_pixel(pixel_height, resized_img.height)

            #prepare panel location
            horizontal = (i % self.column) * (ppt_width_in_pixel / self.column)
            vertical = (i % self.img_iter // self.column) *ppt_height_in_pixel / self.row

            #prepare image location based on panel location
            vertical_position = self.apply_vertical_margin(i, self.row, self.column, vertical, margin_height)
            horizontal_position = horizontal + margin_width

            #add image to a panel
            with io.BytesIO() as output:
                #resized_img.save(output, format="GIF")
                quality_val = 100
                resized_img.save(output,format = "GIF",quality=quality_val)
                image_slide.shapes.add_picture(output, horizontal_position*emus_per_px, vertical_position*emus_per_px)

        self.ppt_component.draw_rectangle(image_slide,self.ppt_width,self.ppt_height)

        return prs

    def apply_vertical_margin(self, index, row, column, vertical, margin_height):
        # apply margin
        iter = row * column
        if 0 <= index % iter and index % iter < column:
            vertical_position = vertical
        elif column * (row - 1) <= index % iter and index % iter < iter:
            vertical_position = vertical + margin_height * 2
        else:
            vertical_position = vertical + margin_height
        return vertical_position

class LoadingConfig(object):
    def __init__(self, dict):
        vars(self).update(dict)

class ConfigObject:
    def __init__(self):
        self.input_path = None
        self.output_path = None

class SlideComponents:
    def textbox(self, image_slide,cellnumber,ppt_width,ppt_height):

        central_box = image_slide.shapes.add_textbox(Inches(ppt_width / 2 - 0.65),
                                                     Inches(ppt_height / 2 - 0.6), Inches(1), Inches(1))
        central_label = central_box.text_frame.add_paragraph()
        central_label.text = "Cell " + cellnumber
        central_label.font.size = Pt(30)


    def draw_rectangle(self, image_slide,ppt_width,ppt_height):
        tx_width = 4
        tx_height = 1
        tx_top = Inches((ppt_height - tx_height) / 2)
        tx_left = Inches((ppt_width - tx_width) / 2)
        rect0 = image_slide.shapes.add_shape(  # shapeオブジェクト➀を追加
            MSO_SHAPE.ROUNDED_RECTANGLE,  # 図形の種類を[丸角四角形]に指定
            tx_left, tx_top,  # 挿入位置の指定　左上の座標の指定
            Inches(tx_width), Inches(tx_height))  # 挿入図形の幅と高さの指定

        rect0.fill.solid()  # shapeオブジェクト➀を単色で塗り潰す
        rect0.fill.fore_color.rgb = RGBColor(250, 100, 100)  # RGB指定で色を指定

        pg = rect0.text_frame.paragraphs[0]  # shapeオブジェクト➀のTextFrameの取得
        pg.text = 'ROUNDED_RECTANGLE'  # TextFrameにテキストを設定
        pg.font.size = Pt(10)  # テキストの文字サイズを10ポイントとする

